"""
Generate robocasa_slides.pptx — clean scientific / conference style.
Run: conda run -n robocasa python make_slides_pptx.py
"""

import json, os, re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette (minimal, scientific) ────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
DGRAY   = RGBColor(0x26, 0x26, 0x26)   # body text
MGRAY   = RGBColor(0x59, 0x59, 0x59)   # secondary text
LGRAY   = RGBColor(0xD0, 0xD0, 0xD0)   # borders / rules
BLUE    = RGBColor(0x1F, 0x4E, 0x79)   # heading / accent
LBLUE   = RGBColor(0x2E, 0x75, 0xB6)   # sub-accent
ORANGE  = RGBColor(0xC5, 0x50, 0x00)
GREEN   = RGBColor(0x37, 0x5E, 0x23)
PURPLE  = RGBColor(0x40, 0x28, 0x68)
RED     = RGBColor(0xC0, 0x00, 0x00)
LBLUE_F = RGBColor(0xD6, 0xE4, 0xF0)   # light fill for header rows

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

OUTPUT = Path("/Users/hossein/Documents/code/robocasa/output")


# ── low-level helpers ─────────────────────────────────────────────────────────

def slide():
    sl = prs.slides.add_slide(blank)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = WHITE
    return sl


def tb(sl, text, x, y, w, h, size=14, bold=False, color=DGRAY,
       align=PP_ALIGN.LEFT, italic=False):
    shape = sl.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    r.font.name = "Calibri"
    return shape


def hline(sl, x, y, w, color=LGRAY, thickness_pt=0.75):
    from pptx.util import Pt
    ln = sl.shapes.add_shape(1, x, y, w, Pt(1))
    ln.fill.background()
    ln.line.color.rgb = color
    ln.line.width = Pt(thickness_pt)


def rect(sl, x, y, w, h, fill=None, line_color=None, line_pt=0.75):
    s = sl.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color; s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s


def slide_heading(sl, title, subtitle=None):
    tb(sl, title, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.55),
       size=24, bold=True, color=BLUE)
    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), color=BLUE, thickness_pt=1.5)
    if subtitle:
        tb(sl, subtitle, Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.38),
           size=12, italic=True, color=MGRAY)
    return Inches(1.05) if subtitle else Inches(0.9)


def bullet_tb(sl, items, x, y, w, h, size=13, indent="  "):
    shape = sl.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.name = "Calibri"
        r.font.color.rgb = DGRAY
    return shape


def simple_table(sl, headers, rows, x, y, w, h, col_widths=None,
                 header_fill=LBLUE_F, hfont=12, bfont=11):
    nc = len(headers)
    nr = len(rows) + 1
    tbl_shape = sl.shapes.add_table(nr, nc, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    def _cell(r, c, text, bold=False, bg=None, fg=DGRAY, sz=bfont, align=PP_ALIGN.LEFT):
        cell = tbl.cell(r, c)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = fg
        if bg:
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
        else:
            cell.fill.background()

    for j, hdr in enumerate(headers):
        _cell(0, j, hdr, bold=True, bg=header_fill, fg=BLUE, sz=hfont)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            is_best = isinstance(val, tuple)
            text  = val[0] if is_best else val
            color = val[1] if is_best else DGRAY
            alt = RGBColor(0xF2,0xF2,0xF2) if i % 2 == 1 else None
            _cell(i + 1, j, text, fg=color, bg=alt)


def img(sl, path, x, y, w, h=None):
    if h:
        sl.shapes.add_picture(str(path), x, y, w, h)
    else:
        sl.shapes.add_picture(str(path), x, y, w)


# ── contact-sheet helpers ─────────────────────────────────────────────────────

def parse_folder(folder_name):
    """Return (task_name, instruction_text) from a folder like Task__slug_text."""
    parts = folder_name.split("__", 1)
    task = parts[0]
    slug = parts[1] if len(parts) > 1 else ""
    instr = slug.replace("_", " ").strip()
    if instr and not instr.endswith("."):
        instr = instr[0].upper() + instr[1:]
    return task, instr


def full_instruction(task_name):
    """Try to get full instruction from episodes.jsonl."""
    from robocasa.utils.dataset_registry_utils import get_ds_path
    for split in ("pretrain", "target"):
        path = get_ds_path(task_name, source="human", split=split)
        if path and os.path.exists(path):
            ep_file = Path(path) / "meta" / "episodes.jsonl"
            if ep_file.exists():
                with open(ep_file) as f:
                    ep = json.loads(f.readline())
                    tasks = ep.get("tasks", [])
                    if tasks:
                        return tasks[0]
    return None


def get_horizon(task_name):
    try:
        from robocasa.utils.dataset_registry import (
            ATOMIC_TASK_DATASETS, COMPOSITE_TASK_DATASETS)
        d = ATOMIC_TASK_DATASETS.get(task_name) or COMPOSITE_TASK_DATASETS.get(task_name)
        return d["horizon"] if d else None
    except Exception:
        return None


def task_type(task_name):
    try:
        from robocasa.utils.dataset_registry import ATOMIC_TASK_DATASETS
        return "Atomic" if task_name in ATOMIC_TASK_DATASETS else "Composite"
    except Exception:
        return ""


def find_contact_sheet_any(task: str):
    """Search all output subdirs for a contact sheet for the given task."""
    search_dirs = [
        OUTPUT / "examples_pure_manipulation" / "atomic",
        OUTPUT / "examples_pure_manipulation" / "composite",
        OUTPUT / "examples_navigation" / "atomic",
        OUTPUT / "examples_navigation" / "composite",
        OUTPUT / "longest" / "atomic",
        OUTPUT / "longest" / "composite",
    ]
    for d in search_dirs:
        if not d.exists():
            continue
        for folder in d.iterdir():
            if folder.name.startswith(task + "__"):
                cs = folder / "contact_sheet.png"
                if cs.exists():
                    return cs
    return None


def task_grid(sl, items, top_y, cols=2, img_w_in=5.9):
    """
    items = list of (contact_sheet_path, task_name, instruction, horizon)
    Lays them out in a grid of `cols` columns starting at top_y.
    img_w_in: width of each contact sheet in inches (height auto from 4:1 ratio).
    """
    img_h_in = img_w_in / 4.0
    label_h  = Inches(0.80)   # task name (0.28) + instruction (0.50) + gap
    cell_h   = Inches(img_h_in) + label_h + Inches(0.15)
    margin   = Inches(0.45)
    gap      = Inches(0.25)

    total_w  = Inches(13.33)
    cell_w_i = (total_w - 2 * margin - (cols - 1) * gap) / cols

    for idx, (cs_path, task, instr, horizon) in enumerate(items):
        col_i = idx % cols
        row_i = idx // cols
        x = margin + col_i * (cell_w_i + gap)
        y = top_y + row_i * (cell_h + Inches(0.2))

        # thin border box
        rect(sl, x, y, cell_w_i, cell_h,
             line_color=LGRAY, line_pt=0.5)

        # image
        img(sl, cs_path, x + Inches(0.05), y + Inches(0.05),
            cell_w_i - Inches(0.1), Inches(img_h_in))

        # task name
        ty = Inches(img_h_in) + y + Inches(0.1)
        h_info = f"  [{task_type(task)} · {horizon} steps]" if horizon else ""
        tb(sl, task + h_info,
           x + Inches(0.08), ty, cell_w_i - Inches(0.12), Inches(0.28),
           size=10, bold=True, color=BLUE)

        # full instruction — no truncation
        instr_text = full_instruction(task) or instr
        tb(sl, instr_text,
           x + Inches(0.08), ty + Inches(0.28),
           cell_w_i - Inches(0.12), Inches(0.50),
           size=9, color=MGRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
tb(sl, "RoboCasa365",
   Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.1),
   size=48, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
hline(sl, Inches(2.5), Inches(2.55), Inches(8.3), BLUE, 1.5)
tb(sl, "A Large-Scale Benchmark for Long-Horizon Robot Learning in Kitchen Environments",
   Inches(1.2), Inches(2.65), Inches(10.9), Inches(0.7),
   size=17, color=DGRAY, align=PP_ALIGN.CENTER)

for i, (val, label) in enumerate([
    ("365 Tasks", "65 atomic + 300 composite"),
    ("2,290+ Hours", "demo data"),
    ("2,500+ Scenes", "diverse kitchen layouts"),
    ("ICLR 2026", "arxiv:2603.04356"),
]):
    bx = Inches(0.5) + i * Inches(3.2)
    rect(sl, bx, Inches(3.55), Inches(2.9), Inches(1.35), line_color=LBLUE, line_pt=1.0)
    tb(sl, val,   bx, Inches(3.65), Inches(2.9), Inches(0.6),
       size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    tb(sl, label, bx, Inches(4.28), Inches(2.9), Inches(0.5),
       size=12, color=MGRAY, align=PP_ALIGN.CENTER)

tb(sl, "robocasa.ai   ·   Hossein Souri et al., ICLR 2026",
   Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
   size=11, color=MGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Overview
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "RoboCasa365 — Overview")

tb(sl,
   "RoboCasa365 is a large-scale simulation benchmark for training and evaluating generalist robot "
   "policies in kitchen environments. The robot must combine mobile navigation, dexterous manipulation, "
   "and long-horizon task planning across 365 diverse tasks.",
   Inches(0.5), top, Inches(12.3), Inches(0.7), size=14, color=DGRAY)

cols_data = [
    ("365", "Total Tasks\n(65 atomic + 300 composite)"),
    ("60",  "Activity Categories"),
    ("2,500+", "Unique Kitchen Scenes"),
    ("2,290+", "Hours of Demo Data"),
    ("2,200+", "3D Kitchen Objects"),
    ("20 Hz", "Control Frequency"),
]
for i, (val, lbl) in enumerate(cols_data):
    bx = Inches(0.5) + i * Inches(2.15)
    rect(sl, bx, top + Inches(0.75), Inches(2.0), Inches(1.2),
         fill=LBLUE_F, line_color=LBLUE, line_pt=0.5)
    tb(sl, val, bx, top + Inches(0.8), Inches(2.0), Inches(0.55),
       size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    tb(sl, lbl, bx, top + Inches(1.32), Inches(2.0), Inches(0.55),
       size=10, color=MGRAY, align=PP_ALIGN.CENTER)

bullet_tb(sl, [
    "Robot:          Franka Panda arm on Omron mobile base (PandaOmron)",
    "Simulator:      MuJoCo via RoboSuite",
    "Pretrain data:  100 human demos × 300 tasks = ~30,000 episodes",
    "Target data:    500 human demos × 50 tasks = ~25,000 episodes",
    "Synthetic data: ~10,000 MimicGen demos × 60 atomic tasks = 600,000+ episodes",
    "Evaluation:     3 protocols — Multi-Task, Foundation Model, Lifelong Learning",
], Inches(0.5), top + Inches(2.15), Inches(12.3), Inches(2.8), size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Task Taxonomy
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Task Taxonomy: Atomic vs. Composite Tasks")

# Atomic box
rect(sl, Inches(0.5), top, Inches(6.0), Inches(5.6),
     line_color=LBLUE, line_pt=1.0)
tb(sl, "Atomic Tasks (65 total)",
   Inches(0.65), top + Inches(0.1), Inches(5.7), Inches(0.4),
   size=15, bold=True, color=BLUE)
hline(sl, Inches(0.65), top + Inches(0.5), Inches(5.7), LBLUE, 0.5)
bullet_tb(sl, [
    "Single-skill, single-step tasks across 14 fixture types",
    "",
    "• Navigation:       NavigateKitchen (1 task)",
    "• Pick & Place:     20 variants (counter ↔ fridge, sink, oven, microwave…)",
    "• Door / Drawer:    Open/Close for 6 appliance types",
    "• Appliance ctrl:   Stove knobs, faucets, blender, toaster…",
    "• Fine manipulation: Adjust temp, turn spout, slide racks",
    "",
    "Horizon:  200 – 700 steps (10 – 35 s @ 20 Hz)",
    "Demos:    ~100 human demos/task · 6,500 total episodes",
    "MimicGen: 60/65 tasks → ~10,000 synthetic demos each",
], Inches(0.65), top + Inches(0.6), Inches(5.7), Inches(4.9), size=12)

# Composite box
rect(sl, Inches(6.83), top, Inches(6.0), Inches(5.6),
     line_color=ORANGE, line_pt=1.0)
tb(sl, "Composite Tasks (300 total — 60 categories)",
   Inches(6.98), top + Inches(0.1), Inches(5.7), Inches(0.4),
   size=15, bold=True, color=ORANGE)
hline(sl, Inches(6.98), top + Inches(0.5), Inches(5.7), ORANGE, 0.5)
bullet_tb(sl, [
    "Multi-step long-horizon tasks — 2 to 10+ sequential sub-goals",
    "",
    "• Require planning over a sequence of manipulation skills",
    "• Cover 60 activity categories (cooking, cleaning, organizing…)",
    "• Example: MakeBananaMilkshake =",
    "    open fridge → pick banana → place in blender",
    "    → close lid → turn on blender → …",
    "",
    "Horizon:  400 – 4,800 steps (20 – 240 s @ 20 Hz)",
    "Demos:    ~100 human demos/task · 31,922 total episodes",
    "Target:   50 tasks held-out for evaluation (16 composite seen,",
    "          16 composite unseen)",
], Inches(6.98), top + Inches(0.6), Inches(5.7), Inches(4.9), size=12)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Navigation
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Navigation & Mobile Manipulation",
                    subtitle="Classification by measured robot base displacement from episode_000000.parquet · threshold > 0.5 m")

tb(sl,
   "Navigation is determined empirically: max base displacement (observation.state[0:2]) from the "
   "start position across the demo episode. Tasks with displacement > 0.5 m require mobile navigation.",
   Inches(0.5), top, Inches(12.3), Inches(0.50), size=12, color=DGRAY)

simple_table(sl,
    ["Category", "Confirmed Navigation (> 0.5 m)", "Pure Manipulation (≤ 0.5 m)"],
    [
        ["Atomic  (65 tasks)",    "1  (NavigateKitchen only)",           "64 / 65  (98.5%)"],
        ["Composite  (300 tasks)", "≥ 13 confirmed from measured sample", "≥ 287 / 300  (≥ 95.7%)"],
        ["Total  (365 tasks)",    ("≥ 14 confirmed", BLUE),              ("≥ 351 / 365  (≥ 96.2%)", GREEN)],
    ],
    Inches(0.5), top + Inches(0.60), Inches(12.3), Inches(1.05),
    col_widths=[Inches(2.0), Inches(4.0), Inches(6.3)])

tb(sl, "All 14 confirmed navigation tasks  (sorted by measured base displacement):",
   Inches(0.5), top + Inches(1.80), Inches(12.3), Inches(0.28),
   size=12, bold=True, color=DGRAY)

simple_table(sl,
    ["Task", "Type", "Horizon", "Displacement", "Full language instruction (episode 0)"],
    [
        ["SearingMeat",           "Composite", "2,900  (145 s)", "1.02 m",
         full_instruction("SearingMeat") or "Grab the pan from the cabinet and place it on the rear right burner…"],
        ["PlaceDishesBySink",     "Composite", "1,600  (80 s)",  "1.21 m",
         full_instruction("PlaceDishesBySink") or "Pick the cup and bowl from the counter and place them by the sink"],
        ["ReturnWashingSupplies", "Composite", "3,400  (170 s)", "1.49 m",
         full_instruction("ReturnWashingSupplies") or "Pick up the bar soap and dish brush from the sink and place them back"],
        ["ArrangeUtensilsByType", "Composite", "3,400  (170 s)", "1.85 m",
         full_instruction("ArrangeUtensilsByType") or "Organize the utensils by type to different locations"],
        ["BeverageSorting",       "Composite", "3,500  (175 s)", "2.40 m",
         full_instruction("BeverageSorting") or "Sort all alcoholic drinks to one cabinet and non-alcoholic to another"],
        ["NavigateKitchen",       "Atomic",    "300  (15 s)",    "2.44 m",
         full_instruction("NavigateKitchen") or "Navigate to the fridge"],
        ["ServeTea",              "Composite", "900  (45 s)",    "2.94 m",
         full_instruction("ServeTea") or "Pick the cup of tea, navigate to the dining table and place it on the saucer"],
        ["HotDogSetup",           "Composite", "2,800  (140 s)", "2.94 m",
         full_instruction("HotDogSetup") or "Move the hot dog bun to the plate on the dining table"],
        ["RecycleSodaCans",       "Composite", "4,500  (225 s)", "3.05 m",
         full_instruction("RecycleSodaCans") or "Four empty soda cans are scattered around the kitchen, gather all the cans"],
        ["PrepareCocktailStation","Composite", "3,500  (175 s)", "3.41 m",
         full_instruction("PrepareCocktailStation") or "There is a bowl on the dining counter, grab a lemon wedge from the fridge"],
        ["DivideBuffetTrays",     "Composite", "4,800  (240 s)", "3.70 m",
         full_instruction("DivideBuffetTrays") or "Gather the lemon wedge and cucumber from the fridge and place them on the trays"],
        ["SetBowlsForSoup",       "Composite", "3,600  (180 s)", "4.52 m",
         full_instruction("SetBowlsForSoup") or "Move the bowls from the cabinet to the plates on the dining table"],
        ["PackFoodByTemp",        "Composite", "3,300  (165 s)", "4.80 m",
         full_instruction("PackFoodByTemp") or "Place the cold items from the fridge in one tupperware and warm items in another"],
        ["PrepareDrinkStation",   "Composite", "3,900  (195 s)", "6.93 m",
         full_instruction("PrepareDrinkStation") or "There is a tray on the dining counter, grab the cup, mug and pitcher"],
    ],
    Inches(0.5), top + Inches(2.14), Inches(12.3), Inches(4.05),
    col_widths=[Inches(2.1), Inches(1.0), Inches(1.3), Inches(1.0), Inches(6.9)],
    hfont=10, bfont=8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5a — Target Evaluation Tasks (overview)
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Target Evaluation Tasks — Overview",
                    subtitle="50 held-out tasks used for benchmarking across 3 splits · 10 unseen kitchen scenes")

# summary boxes
for i, (val, label, sub, col) in enumerate([
    ("50",  "Total target tasks", "18 atomic + 32 composite", BLUE),
    ("49",  "Pure manipulation",  "98% of target tasks",      GREEN),
    ("1",   "Navigation task",    "NavigateKitchen (atomic)", ORANGE),
    ("10",  "Held-out scenes",    "Disjoint from training",   PURPLE),
]):
    bx = Inches(0.5) + i * Inches(3.2)
    rect(sl, bx, top, Inches(2.9), Inches(1.25), fill=LBLUE_F, line_color=LBLUE, line_pt=0.5)
    tb(sl, val,   bx, top + Inches(0.05), Inches(2.9), Inches(0.55),
       size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(sl, label, bx, top + Inches(0.6),  Inches(2.9), Inches(0.32),
       size=11, bold=True, color=DGRAY, align=PP_ALIGN.CENTER)
    tb(sl, sub,   bx, top + Inches(0.9),  Inches(2.9), Inches(0.28),
       size=9, color=MGRAY, align=PP_ALIGN.CENTER)

# three-column split table
simple_table(sl,
    ["Split", "Tasks", "Task type", "Horizon range", "Navigation"],
    [
        ["Atomic-seen",      "18", "Atomic",    "300–700 steps  (15–35 s)",    "1 task (NavigateKitchen)"],
        ["Composite-seen",   "16", "Composite", "800–2,900 steps  (40–145 s)", "None"],
        ["Composite-unseen", "16", "Composite", "800–2,900 steps  (40–135 s)", "None"],
        ["Total",            ("50", BLUE), "—", "300–2,900 steps", ("1 / 50  (2%)", ORANGE)],
    ],
    Inches(0.5), top + Inches(1.45), Inches(12.3), Inches(1.7),
    col_widths=[Inches(1.9), Inches(0.8), Inches(1.3), Inches(3.3), Inches(5.0)])

tb(sl, "Atomic-seen vs. Composite-seen vs. Composite-unseen:",
   Inches(0.5), top + Inches(3.35), Inches(12.3), Inches(0.3),
   size=13, bold=True, color=DGRAY)
bullet_tb(sl, [
    "Atomic-seen:        18 atomic tasks the model was trained on — tests atomic skill retention",
    "Composite-seen:     16 composite tasks the model was trained on — tests multi-step execution on held-out scenes",
    "Composite-unseen:   16 composite tasks NEVER seen during training — tests compositional generalization",
], Inches(0.5), top + Inches(3.7), Inches(12.3), Inches(1.5), size=12)

rect(sl, Inches(0.5), top + Inches(5.35), Inches(12.3), Inches(0.62),
     fill=RGBColor(0xFF, 0xF2, 0xCC), line_color=ORANGE, line_pt=0.75)
tb(sl, "Composite-unseen is the hardest split: tasks are novel compositions of learned skills "
   "evaluated in unseen kitchen scenes. Best model (GR00T N1.5) achieves only 4.4%.",
   Inches(0.65), top + Inches(5.4), Inches(12.0), Inches(0.54), size=12, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5b — Target Evaluation Tasks (full list)
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Target Evaluation Tasks — Full List")

atomic_seen = [
    ("CloseBlenderLid","600"),("CloseFridge","600"),("CloseToasterOvenDoor","300"),
    ("CoffeeSetupMug","400"),("NavigateKitchen","300"),("OpenCabinet","700"),
    ("OpenDrawer","500"),("OpenStandMixerHead","300"),("PickPlaceCounterToCabinet","500"),
    ("PickPlaceCounterToStove","400"),("PickPlaceDrawerToCounter","500"),
    ("PickPlaceSinkToCounter","600"),("PickPlaceToasterToCounter","400"),
    ("SlideDishwasherRack","300"),("TurnOffStove","500"),
    ("TurnOnElectricKettle","300"),("TurnOnMicrowave","300"),("TurnOnSinkFaucet","400"),
]
composite_seen = [
    ("DeliverStraw","1700"),("GetToastedBread","2000"),("KettleBoiling","1000"),
    ("LoadDishwasher","1200"),("PackIdenticalLunches","2600"),("PreSoakPan","1600"),
    ("PrepareCoffee","1200"),("RinseSinkBasin","900"),("ScrubCuttingBoard","800"),
    ("SearingMeat","2900"),("SetUpCuttingStation","1600"),("StackBowlsCabinet","1400"),
    ("SteamInMicrowave","1400"),("StirVegetables","1600"),("StoreLeftoversInBowl","1700"),
    ("WashLettuce","1100"),
]
composite_unseen = [
    ("ArrangeBreadBasket","2900"),("ArrangeTea","1500"),("BreadSelection","1300"),
    ("CategorizeCondiments","1100"),("CuttingToolSelection","800"),("GarnishPancake","1800"),
    ("GatherTableware","1500"),("HeatKebabSandwich","1800"),("MakeIceLemonade","2000"),
    ("PanTransfer","1200"),("PortionHotDogs","1500"),("RecycleBottlesByType","1900"),
    ("SeparateFreezerRack","1600"),("WaffleReheat","2700"),("WashFruitColander","2100"),
    ("WeighIngredients","2000"),
]

col_w = Inches(4.05)
col_gap = Inches(0.07)

for col_i, (title, col_color, tasks) in enumerate([
    ("Atomic-seen  (18 tasks)",       BLUE,   atomic_seen),
    ("Composite-seen  (16 tasks)",    LBLUE,  composite_seen),
    ("Composite-unseen  (16 tasks)",  ORANGE, composite_unseen),
]):
    cx = Inches(0.5) + col_i * (col_w + col_gap)
    # header
    rect(sl, cx, top, col_w, Inches(0.36), fill=col_color)
    tb(sl, title, cx + Inches(0.08), top + Inches(0.04), col_w - Inches(0.1), Inches(0.3),
       size=12, bold=True, color=WHITE)
    # rows
    row_h = Inches(0.34)
    for row_i, (task, horizon) in enumerate(tasks):
        ry = top + Inches(0.4) + row_i * row_h
        fill = RGBColor(0xF2,0xF2,0xF2) if row_i % 2 == 0 else WHITE
        rect(sl, cx, ry, col_w, row_h, fill=fill, line_color=LGRAY, line_pt=0.3)
        nav_mark = "  ★NAV" if task == "NavigateKitchen" else ""
        tb(sl, task + nav_mark,
           cx + Inches(0.08), ry + Inches(0.05), Inches(2.9), Inches(0.26),
           size=9.5, color=ORANGE if nav_mark else DGRAY, bold=bool(nav_mark))
        tb(sl, f"{int(horizon)//20}s  ({horizon} steps)",
           cx + Inches(3.0), ry + Inches(0.05), Inches(1.0), Inches(0.26),
           size=8.5, color=MGRAY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5c — Training & Evaluation Data
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Training & Evaluation Data",
                    subtitle="Dataset splits used for pretraining and held-out evaluation")

# ── Training table ─────────────────────────────────────────────────────────────
tb(sl, "Training Data  (Pretrain Split)",
   Inches(0.5), top, Inches(12.3), Inches(0.32),
   size=13, bold=True, color=BLUE)
simple_table(sl,
    ["Source", "Split", "# Tasks", "Demos / Task", "Total Demos", "Total Hours"],
    [
        ["Human teleoperation",  "Pretrain", "300", "~100",    "~30,000",  "482 hrs"],
        ["MimicGen (synthetic)", "Pretrain", "60",  "~10,000", "~600,000", "1,615 hrs"],
        [("Total", BLUE), "—", ("360", BLUE), "—", ("~630,000", BLUE), ("~2,097 hrs", BLUE)],
    ],
    Inches(0.5), top + Inches(0.38), Inches(12.3), Inches(1.55),
    col_widths=[Inches(2.5), Inches(1.0), Inches(1.0), Inches(1.5), Inches(1.7), Inches(4.6)])

# ── Evaluation table ───────────────────────────────────────────────────────────
tb(sl, "Evaluation Data  (Target Split — RoboCasa365, ICLR 2026)",
   Inches(0.5), top + Inches(2.1), Inches(12.3), Inches(0.32),
   size=13, bold=True, color=BLUE)
simple_table(sl,
    ["Source", "Split", "# Tasks", "Subset", "Demos / Task", "Total Demos"],
    [
        ["Human teleoperation", "Target", "18", "Atomic-seen",      "~500", "~9,000"],
        ["Human teleoperation", "Target", "16", "Composite-seen",   "~500", "~8,000"],
        ["Human teleoperation", "Target", "16", "Composite-unseen", "~500", "~8,000"],
        [("Total", BLUE), "—", ("50", BLUE), "—", "~500", ("~25,000", BLUE)],
    ],
    Inches(0.5), top + Inches(2.48), Inches(12.3), Inches(1.85),
    col_widths=[Inches(2.5), Inches(1.0), Inches(1.0), Inches(2.0), Inches(1.5), Inches(4.3)])

# ── note ───────────────────────────────────────────────────────────────────────
rect(sl, Inches(0.5), top + Inches(4.5), Inches(12.3), Inches(0.75),
     fill=LBLUE_F, line_color=LBLUE, line_pt=0.75)
tb(sl,
   "Note: The original RoboCasa (RSS 2024, arxiv:2406.02523) used 100 evaluation tasks. "
   "RoboCasa365 (ICLR 2026, arxiv:2603.04356) defines a new 50-task target split with "
   "held-out kitchen scenes and 500 demos/task for more rigorous evaluation.",
   Inches(0.65), top + Inches(4.55), Inches(12.0), Inches(0.65),
   size=11, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Task Durations
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Task Durations")

# left: distribution table
tb(sl, "Composite task duration distribution",
   Inches(0.5), top, Inches(6.0), Inches(0.32), size=13, bold=True, color=DGRAY)
simple_table(sl,
    ["Duration Range", "Count", "Percentage"],
    [
        ["< 30 s  (very short)",  "65",  "21%"],
        ["30 – 60 s  (short)",    "141", "46%"],
        ["60 – 100 s  (medium)",  "81",  "27%"],
        ["> 100 s  (long)",       "17",  "6%"],
    ],
    Inches(0.5), top + Inches(0.38), Inches(5.8), Inches(1.55))

tb(sl, "Summary statistics:",
   Inches(0.5), top + Inches(2.08), Inches(5.8), Inches(0.3), size=12, bold=True, color=DGRAY)
simple_table(sl,
    ["Min", "Mean", "Median", "Max"],
    [["10 s", "52 s", "46 s", "240 s"]],
    Inches(0.5), top + Inches(2.42), Inches(5.8), Inches(0.75))

# right: top 10 longest
tb(sl, "10 longest tasks (by max horizon)",
   Inches(6.8), top, Inches(6.0), Inches(0.32), size=13, bold=True, color=DGRAY)
simple_table(sl,
    ["Rank", "Task", "Steps", "Duration"],
    [
        ["1",  "DivideBuffetTrays",     "4,800", "240 s"],
        ["2",  "RecycleSodaCans",       "4,500", "225 s"],
        ["3",  "PrepareDrinkStation",   "3,900", "195 s"],
        ["4",  "SetBowlsForSoup",       "3,600", "180 s"],
        ["4",  "SpicyMarinade",         "3,600", "180 s"],
        ["6",  "BeverageSorting",       "3,500", "175 s"],
        ["6",  "PrepareCocktailStation","3,500", "175 s"],
        ["8",  "ArrangeUtensilsByType", "3,400", "170 s"],
        ["8",  "ReturnWashingSupplies", "3,400", "170 s"],
        ["10", "ClearSink / PackFoodByTemp","3,300","165 s"],
    ],
    Inches(6.8), top + Inches(0.38), Inches(6.0), Inches(3.6),
    col_widths=[Inches(0.5), Inches(2.7), Inches(1.1), Inches(1.7)])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Benchmark Comparison
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Benchmark Comparison")

simple_table(sl,
    ["Feature", "LIBERO", "LIBERO-Long", "MetaWorld", "Calvin", "RoboCasa365"],
    [
        ["Total tasks",               "50",       "10",         "50",      "34",     ("365", BLUE)],
        ["Long-horizon tasks",         "Partial",  "4–6 steps",  "No",      "Chained",("300 (2–10+ steps)", BLUE)],
        ["Typical episode length",     "15–30 s",  "60–120 s",   "5–15 s",  "30–60 s",("10–240 s (avg 52 s)", BLUE)],
        ["Navigation / mobile base",   "No",       "No",         "No",      "Yes",    ("Yes — full mobile base", BLUE)],
        ["Scene diversity",            "~10",      "~10",        "Single",  "1",      ("2,500+", BLUE)],
        ["Total demo data",            "~50 hrs",  "~10 hrs",    "Scripted","~24 hrs",("2,290+ hrs", BLUE)],
        ["Activity categories",        "4",        "1",          "12",      "4",      ("60", BLUE)],
        ["Held-out generalization",    "Limited",  "Limited",    "No",      "Partial",("Yes — disjoint target scenes", BLUE)],
        ["Lifelong learning protocol", "Yes",      "No",         "No",      "No",     ("Yes — 4-phase", BLUE)],
    ],
    Inches(0.5), top, Inches(12.3), Inches(5.5),
    col_widths=[Inches(2.5), Inches(1.3), Inches(1.55), Inches(1.3), Inches(1.2), Inches(4.45)])

tb(sl, "* LIBERO numbers from published paper.  RoboCasa365 from ICLR 2026 paper.",
   Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3), size=9, color=MGRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Why Long-Horizon
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Why RoboCasa365 for Long-Horizon Reasoning?")

reasons = [
    ("1. Sequential Decision Chains",
     "Tasks require 2–10+ sub-goals in the correct order. "
     "Failure at any step cascades — no partial credit from atomic success."),
    ("2. Semantic Reasoning Required",
     "Tasks like LoadFridgeFifo, BeverageSorting, FilterMicrowavableItem require "
     "understanding categories, not only mimicking motor patterns."),
    ("3. Planning Under Partial Observability",
     "Fridge compartments, drawers, and cabinets are hidden. "
     "The robot must plan retrieval sequences without full state visibility."),
    ("4. Temporal Depth",
     "Tasks span 200–4,800 control steps per episode at 20 Hz. "
     "This directly tests long-horizon credit assignment in policy learning."),
    ("5. Compositional Generalization",
     "Composite-Unseen tasks test whether learned skills compose to novel task sequences — "
     "a key measure of genuine planning vs. memorization."),
    ("6. Realistic Action Space",
     "Appliance control (knobs, lids, buttons) + navigation + dexterous grasping = "
     "the full action vocabulary of household manipulation robots."),
]

cw = Inches(6.0)
ch = Inches(1.2)
for idx, (title, body) in enumerate(reasons):
    cx = Inches(0.5) + (idx % 2) * (cw + Inches(0.33))
    cy = top + (idx // 2) * (ch + Inches(0.1))
    rect(sl, cx, cy, cw, ch, line_color=LGRAY, line_pt=0.5)
    rect(sl, cx, cy, Pt(4), ch, fill=BLUE)
    tb(sl, title, cx + Inches(0.1), cy + Inches(0.08), cw - Inches(0.15), Inches(0.32),
       size=12, bold=True, color=BLUE)
    tb(sl, body, cx + Inches(0.1), cy + Inches(0.4), cw - Inches(0.15), Inches(0.75),
       size=11, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Results
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Benchmark Results")

tb(sl, "Multi-Task Learning  (train on 300 tasks, 482 hrs human demos)",
   Inches(0.5), top, Inches(12.3), Inches(0.3), size=13, bold=True, color=DGRAY)
simple_table(sl,
    ["Task Split", "Diffusion Policy", "π₀", "π₀.₅", "GR00T N1.5"],
    [
        ["Atomic-Seen",      "15.7%", "34.6%", "39.6%", ("43.0%", GREEN)],
        ["Composite-Seen",   "0.2%",  "6.1%",  "7.1%",  ("9.6%",  GREEN)],
        ["Composite-Unseen", "1.25%", "1.1%",  "1.2%",  ("4.4%",  GREEN)],
        ["Average",          "6.1%",  "14.8%", "16.9%", ("20.0%", GREEN)],
    ],
    Inches(0.5), top + Inches(0.35), Inches(12.3), Inches(1.6),
    col_widths=[Inches(2.2), Inches(2.5), Inches(2.0), Inches(2.0), Inches(3.6)])

tb(sl, "Foundation Model Learning  (pretrain 2,290 hrs → fine-tune 50 target tasks)",
   Inches(0.5), top + Inches(2.1), Inches(12.3), Inches(0.3), size=13, bold=True, color=DGRAY)
simple_table(sl,
    ["Task Split", "Pretrain Only", "Target Only (100%)", "Pretrain + Target (100%)"],
    [
        ["Atomic-Seen",      "41.9%", "60.6%", ("68.5%", GREEN)],
        ["Composite-Seen",   "0.0%",  "35.0%", ("40.6%", GREEN)],
        ["Composite-Unseen", "0.2%",  "33.3%", ("42.1%", GREEN)],
        ["Average",          "15.1%", "43.7%", ("51.1%", GREEN)],
    ],
    Inches(0.5), top + Inches(2.45), Inches(12.3), Inches(1.6),
    col_widths=[Inches(2.2), Inches(2.5), Inches(3.8), Inches(3.8)])

rect(sl, Inches(0.5), top + Inches(4.2), Inches(12.3), Inches(0.75),
     fill=RGBColor(0xFF, 0xF2, 0xCC), line_color=ORANGE, line_pt=0.75)
tb(sl, "Key finding: Even the best model (GR00T N1.5) achieves only 9.6% on composite tasks in "
   "the multi-task setting. These tasks are genuinely challenging and leave significant room for improvement.",
   Inches(0.65), top + Inches(4.25), Inches(12.0), Inches(0.65), size=12, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Three Protocols
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Three Benchmarking Protocols")

protocols = [
    ("1. Multi-Task Learning", BLUE, [
        "Goal: Learn a single policy across all 300 tasks",
        "Training data: 482 hrs human demos, 300 tasks, 100 demos/task",
        "Evaluation: Pretrain kitchen scenes",
        "Measures: Breadth of multi-task generalization",
    ]),
    ("2. Foundation Model Learning", LBLUE, [
        "Goal: Pretrain on large data, fine-tune on 50 target tasks",
        "Pretrain: 482 hrs human + 1,615 hrs MimicGen synthetic",
        "Fine-tune: 193 hrs across 3 splits (atomic-seen, composite-seen, unseen)",
        "Evaluation: Held-out target kitchen scenes",
        "Measures: Transfer learning, data efficiency",
    ]),
    ("3. Lifelong Learning", ORANGE, [
        "Goal: Sequentially learn harder tasks without catastrophic forgetting",
        "Phase 1: Atomic tasks       (~41.5% success)",
        "Phase 2: 2–3 stage tasks   (~24.5%)",
        "Phase 3: 4–5 stage tasks",
        "Phase 4: 6+ stage tasks",
        "Measures: Continual learning, backward transfer (forgetting)",
    ]),
]

cw = Inches(4.0)
for i, (title, col, items) in enumerate(protocols):
    cx = Inches(0.5) + i * (cw + Inches(0.25))
    rect(sl, cx, top, cw, Inches(5.5), line_color=col, line_pt=1.2)
    rect(sl, cx, top, cw, Inches(0.45), fill=col)
    tb(sl, title, cx + Inches(0.1), top + Inches(0.06),
       cw - Inches(0.15), Inches(0.35), size=13, bold=True, color=WHITE)
    bullet_tb(sl, items, cx + Inches(0.1), top + Inches(0.55),
              cw - Inches(0.15), Inches(4.85), size=12)

tb(sl, "Supported models: Diffusion Policy · π₀ / π₀.₅ (OpenPI) · GR00T N1.5",
   Inches(0.5), top + Inches(5.65), Inches(12.3), Inches(0.4),
   size=12, color=MGRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Atomic Task Examples
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Atomic Task Examples",
                    subtitle="Single-skill manipulation tasks — 65 tasks total, horizon 200–700 steps")

# Collect the 5 atomic pure manipulation contact sheets (ordered short to long)
atomic_order = [
    "TurnOnToaster", "TurnOnStove",
    "PickPlaceCounterToStove", "PickPlaceCounterToCabinet", "OpenCabinet",
]
atomic_items = []
for task in atomic_order:
    cs_dir = OUTPUT / "examples_pure_manipulation" / "atomic"
    for d in cs_dir.iterdir():
        if d.name.startswith(task + "__"):
            cs = d / "contact_sheet.png"
            if cs.exists():
                _, instr = parse_folder(d.name)
                h = get_horizon(task)
                atomic_items.append((cs, task, instr, h))
            break

task_grid(sl, atomic_items, top, cols=2, img_w_in=6.0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Pure Manipulation: Short to Medium Composite Tasks
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Pure Manipulation — Short to Medium Composite Tasks",
                    subtitle="Confirmed pure: max base displacement ≤ 0.5 m in demo data · horizon 1,100–3,600 steps")

# All confirmed pure composite tasks from measurement (< 0.01 m displacement)
comp_short = ["WashLettuce", "LoadDishwasher", "PrepareCoffee", "SpicyMarinade"]
items = []
for task in comp_short:
    cs = find_contact_sheet_any(task)
    if cs:
        items.append((cs, task, "", get_horizon(task)))

task_grid(sl, items, top, cols=2, img_w_in=6.0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — 10 Longest Pure Manipulation Tasks
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "10 Longest Pure Manipulation Tasks",
                    subtitle="All confirmed zero base movement (< 0.01 m displacement) · ranked by task horizon")

# Table of all 10 with measured displacement
longest_pure = [
    ("SpicyMarinade",          3600, 0.0002),
    ("ClearSink",              3300, 0.0002),
    ("MicrowaveThawing",       3100, 0.0002),
    ("CerealAndBowl",          2900, 0.0013),
    ("SweetSavoryToastSetup",  2700, 0.0009),
    ("OrganizeCleaningSupplies",2700, 0.0009),
    ("ToastBaguette",          2600, 0.0009),
    ("PrepareToast",           2500, 0.0002),
    ("RestockBowls",           2400, 0.0002),
    ("MakeFruitBowl",          2300, 0.0006),
]
simple_table(sl,
    ["Rank", "Task", "Horizon", "Duration", "Max Base Displacement", "Full language instruction (episode 0)"],
    [
        [str(i+1), task,
         f"{h:,} steps",
         f"{h//20} s",
         f"{d:.4f} m",
         full_instruction(task) or ""]
        for i, (task, h, d) in enumerate(longest_pure)
    ],
    Inches(0.5), top, Inches(12.3), Inches(4.2),
    col_widths=[Inches(0.45), Inches(2.3), Inches(1.1), Inches(0.8), Inches(1.5), Inches(6.15)],
    hfont=11, bfont=10)

# Visual examples for the two tasks that have contact sheets
tb(sl, "Visual examples (contact sheets available):",
   Inches(0.5), top + Inches(4.35), Inches(12.3), Inches(0.28),
   size=11, bold=True, color=DGRAY)
long_pure_items = []
for task, h, _ in longest_pure[:2]:
    cs = find_contact_sheet_any(task)
    if cs:
        long_pure_items.append((cs, task, "", h))
if long_pure_items:
    task_grid(sl, long_pure_items, top + Inches(4.65), cols=2, img_w_in=6.0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Pure Manipulation: From Short to Longest (overview)
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Pure Manipulation — Full Range: Short to Longest",
                    subtitle="All confirmed pure (max base displacement ≤ 0.5 m) · atomic 200–700 steps, composite up to 3,600 steps")

# 5 atomic + 5 pure composite spanning the horizon range
all_10_pure = [
    ("TurnOnToaster",            "atomic"),
    ("TurnOnStove",               "atomic"),
    ("PickPlaceCounterToStove",   "atomic"),
    ("PickPlaceCounterToCabinet", "atomic"),
    ("OpenCabinet",               "atomic"),
    ("WashLettuce",               "composite"),
    ("LoadDishwasher",            "composite"),
    ("PrepareCoffee",             "composite"),
    ("SpicyMarinade",             "composite"),
    ("ClearSink",                 "composite"),
]

simple_table(sl,
    ["#", "Task", "Type", "Horizon", "Duration", "Full language instruction (episode 0)"],
    [
        [str(i+1), task, ttype.capitalize(),
         f"{get_horizon(task):,}" if get_horizon(task) else "—",
         f"{get_horizon(task)//20} s" if get_horizon(task) else "—",
         full_instruction(task) or "—"]
        for i, (task, ttype) in enumerate(all_10_pure)
    ],
    Inches(0.5), top, Inches(12.3), Inches(5.8),
    col_widths=[Inches(0.35), Inches(2.2), Inches(1.1), Inches(0.9),
                Inches(0.9), Inches(6.85)])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Navigation Task Examples (short-horizon)
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Navigation Task Examples — Short & Medium Horizon",
                    subtitle="Confirmed navigation: max base displacement > 0.5 m in demo data")

nav_short = ["NavigateKitchen", "ServeTea", "PlaceDishesBySink", "HotDogSetup"]
nav_items = []
for task in nav_short:
    cs = find_contact_sheet_any(task)
    if cs:
        nav_items.append((cs, task, "", get_horizon(task)))

task_grid(sl, nav_items, top, cols=2, img_w_in=6.0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14b — Navigation Task Examples (long-horizon)
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Navigation Task Examples — Long Horizon",
                    subtitle="All confirmed by measured base displacement > 0.5 m · horizon 2,900–4,800 steps")

nav_long = ["SearingMeat", "ReturnWashingSupplies", "BeverageSorting", "ArrangeUtensilsByType"]
nav_long_items = []
for task in nav_long:
    cs = find_contact_sheet_any(task)
    if cs:
        nav_long_items.append((cs, task, "", get_horizon(task)))

task_grid(sl, nav_long_items, top, cols=2, img_w_in=6.0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14c — Navigation Task Examples (very long horizon)
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Navigation Task Examples — Very Long Horizon",
                    subtitle="Large-scale cross-kitchen tasks · horizon 3,300–4,800 steps · base displacement 3–7 m")

nav_vlong = ["RecycleSodaCans", "DivideBuffetTrays", "SetBowlsForSoup", "PackFoodByTemp",
             "PrepareCocktailStation", "PrepareDrinkStation"]
nav_vlong_items = []
for task in nav_vlong:
    cs = find_contact_sheet_any(task)
    if cs:
        nav_vlong_items.append((cs, task, "", get_horizon(task)))

task_grid(sl, nav_vlong_items, top, cols=2, img_w_in=6.0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — 10 Longest Pure Manipulation Tasks (visual examples)
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "10 Longest Pure Manipulation Tasks — Visual Examples",
                    subtitle="SpicyMarinade & ClearSink contact sheets · all 10 tasks confirmed zero base movement (< 0.01 m)")

# Show the two tasks that have rendered contact sheets
longest_pure_visual = ["SpicyMarinade", "ClearSink"]
items = []
for task in longest_pure_visual:
    cs = find_contact_sheet_any(task)
    if cs:
        items.append((cs, task, "", get_horizon(task)))

task_grid(sl, items, top, cols=2, img_w_in=6.0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Summary
# ═══════════════════════════════════════════════════════════════════════════════
sl = slide()
top = slide_heading(sl, "Summary")

summary_rows = [
    ("Scale",             "365 tasks (65 atomic + 300 composite) · 60 categories · 2,500+ scenes · 2,290+ hours"),
    ("Task Duration",     "10–240 s per episode (avg 52 s) — up to 4,800 control steps"),
    ("Navigation",        "≥ 14 tasks require mobile-base navigation (measured via base displacement > 0.5 m); ≥ 351/365 (≥ 96.2%) are pure manipulation"),
    ("Data",              "675 hrs human + 1,615 hrs MimicGen synthetic = 2,290+ total hours"),
    ("Benchmarking",      "3 protocols: Multi-Task Learning · Foundation Model Learning · Lifelong Learning"),
    ("Best SOTA result",  "GR00T N1.5: 43% atomic-seen, 9.6% composite-seen, 4.4% composite-unseen"),
    ("Why use it",        "Long-horizon, semantically diverse, compositionally challenging — closest to real household robots"),
]

cw = Inches(12.3)
ch = Inches(0.72)
for i, (label, body) in enumerate(summary_rows):
    cy = top + i * (ch + Inches(0.04))
    fill = LBLUE_F if i % 2 == 0 else WHITE
    rect(sl, Inches(0.5), cy, cw, ch, fill=fill, line_color=LGRAY, line_pt=0.4)
    tb(sl, label, Inches(0.65), cy + Inches(0.13), Inches(2.0), Inches(0.48),
       size=12, bold=True, color=BLUE)
    tb(sl, body,  Inches(2.75), cy + Inches(0.13), Inches(9.9), Inches(0.48),
       size=12, color=DGRAY)

tb(sl, "robocasa.ai  ·  arxiv:2603.04356  ·  ICLR 2026",
   Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
   size=11, color=MGRAY, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# GALLERY SLIDES — All Rendered Output Examples
# ═══════════════════════════════════════════════════════════════════════════════
# Confirmed pure manipulation tasks (max base displacement < 0.01 m)
PURE_TASKS = {
    "TurnOnToaster", "TurnOnStove", "PickPlaceCounterToStove",
    "PickPlaceCounterToCabinet", "OpenCabinet",
    "WashLettuce", "LoadDishwasher", "PrepareCoffee", "SpicyMarinade", "ClearSink",
    "MicrowaveThawing", "CerealAndBowl", "SweetSavoryToastSetup",
    "OrganizeCleaningSupplies", "ToastBaguette", "PrepareToast",
    "RestockBowls", "MakeFruitBowl",
}
# Confirmed navigation tasks (base displacement > 0.5 m)
NAV_TASKS = {
    "NavigateKitchen", "ServeTea", "PlaceDishesBySink", "HotDogSetup",
    "SearingMeat", "ReturnWashingSupplies", "BeverageSorting",
    "ArrangeUtensilsByType", "DivideBuffetTrays", "RecycleSodaCans",
    "PackFoodByTemp", "PrepareDrinkStation", "SetBowlsForSoup", "PrepareCocktailStation",
}


def collect_all_contacts(filter_set=None):
    """Collect unique contact sheets across all output dirs. Deduplicates by task name."""
    all_dirs = [
        OUTPUT / "examples_pure_manipulation" / "atomic",
        OUTPUT / "examples_pure_manipulation" / "composite",
        OUTPUT / "examples_navigation" / "atomic",
        OUTPUT / "examples_navigation" / "composite",
        OUTPUT / "longest" / "atomic",
        OUTPUT / "longest" / "composite",
    ]
    seen, items = set(), []
    for d in all_dirs:
        if not d.exists():
            continue
        for folder in sorted(d.iterdir()):
            cs = folder / "contact_sheet.png"
            if not cs.exists():
                continue
            task_name = folder.name.split("__")[0]
            if task_name in seen:
                continue
            if filter_set is not None and task_name not in filter_set:
                continue
            seen.add(task_name)
            instr = full_instruction(task_name) or ""
            items.append((cs, task_name, instr, get_horizon(task_name)))
    return items


def gallery_slides(title, items, cols=2, img_w_in=6.0):
    """Paginate items across gallery slides, 4 per slide (2×2)."""
    per_page = cols * 2
    total_pages = max(1, (len(items) + per_page - 1) // per_page)
    for page in range(total_pages):
        chunk = items[page * per_page:(page + 1) * per_page]
        if not chunk:
            continue
        sl = slide()
        pg_sfx = f"  (Part {page+1} of {total_pages})" if total_pages > 1 else ""
        top = slide_heading(sl, title + pg_sfx)
        task_grid(sl, chunk, top, cols=cols, img_w_in=img_w_in)


# ── Gallery: Pure Manipulation (atomic) ──────────────────────────────────────
atomic_pure = collect_all_contacts(
    filter_set={"TurnOnToaster","TurnOnStove","PickPlaceCounterToStove",
                "PickPlaceCounterToCabinet","OpenCabinet"})
gallery_slides("Gallery — Pure Manipulation: Atomic Tasks", atomic_pure, cols=2)

# ── Gallery: Pure Manipulation (composite) ──────────────────────────────────
comp_pure = collect_all_contacts(
    filter_set=PURE_TASKS - {"TurnOnToaster","TurnOnStove","PickPlaceCounterToStove",
                              "PickPlaceCounterToCabinet","OpenCabinet"})
gallery_slides("Gallery — Pure Manipulation: Composite Tasks", comp_pure, cols=2)

# ── Gallery: Navigation Tasks ────────────────────────────────────────────────
nav_all = collect_all_contacts(filter_set=NAV_TASKS)
gallery_slides("Gallery — Navigation Tasks", nav_all, cols=2)


# ═══════════════════════════════════════════════════════════════════════════════
# APPENDIX — All 365 Task Names Organised by Category
# ═══════════════════════════════════════════════════════════════════════════════

import re as _re

_MD_PATH = Path("/Users/hossein/Documents/code/robocasa/ROBOCASA_BENCHMARK.md")

def _parse_categories():
    """Parse all 60 categories + tasks from ROBOCASA_BENCHMARK.md section 13."""
    text = _MD_PATH.read_text()
    m = _re.search(r'## 13\. Complete Composite Task List.*?(?=\n## 14\.)', text, _re.DOTALL)
    if not m:
        return []
    section = m.group(0)
    cat_hdrs = list(_re.finditer(r'### (\d+)\. (.+?) \((\d+) tasks?\)', section))
    cats = []
    for i, hdr in enumerate(cat_hdrs):
        start = hdr.end()
        end = cat_hdrs[i+1].start() if i+1 < len(cat_hdrs) else len(section)
        chunk = section[start:end]
        tasks = []
        for line in chunk.splitlines():
            parts = [p.strip() for p in line.split('|')]
            if len(parts) < 4:
                continue
            name = parts[1].replace('**', '').replace('*', '').strip()
            if (not name or name in ('Task', '------', '---')
                    or not name[0].isupper() or ' ' in name[:2]):
                continue
            dur  = parts[3].replace('**', '').strip() if len(parts) > 3 else ''
            tgt  = parts[5].strip() if len(parts) > 5 else ''
            tasks.append((name, dur, tgt))
        if tasks:
            cats.append((int(hdr.group(1)), hdr.group(2), tasks))
    return cats

_CATEGORIES = _parse_categories()

# ── Appendix layout constants ─────────────────────────────────────────────────
_COLS       = 3
_APP_MARGIN = Inches(0.36)
_APP_GAP    = Inches(0.16)
_APP_COL_W  = (Inches(13.33) - 2 * _APP_MARGIN - (_COLS - 1) * _APP_GAP) / _COLS
_TASK_H     = Inches(0.193)
_CAT_HDR_H  = Inches(0.280)
_CAT_GAP    = Inches(0.10)
_MAX_COL_H  = Inches(6.22)

# Palette for category headers cycling through 3 accent colours
_CAT_FILLS  = [LBLUE_F, RGBColor(0xED, 0xF7, 0xED), RGBColor(0xFF, 0xF8, 0xE7)]
_CAT_LINES  = [LBLUE,   GREEN,                       ORANGE]
_CAT_TEXTS  = [BLUE,    GREEN,                       ORANGE]


def _col_x(ci):
    return _APP_MARGIN + ci * (_APP_COL_W + _APP_GAP)


def _draw_category(sl, col_i, y, cat_num, cat_name, tasks, page_cat_idx):
    """Draw one category block (header + task rows) at the given column / y."""
    cx = _col_x(col_i)
    ci = page_cat_idx % 3            # cycles colour every 3 categories on the page
    fill  = _CAT_FILLS[ci]
    lc    = _CAT_LINES[ci]
    tc    = _CAT_TEXTS[ci]

    # Header
    rect(sl, cx, y, _APP_COL_W, _CAT_HDR_H, fill=fill, line_color=lc, line_pt=0.5)
    tb(sl, f"{cat_num}. {cat_name}",
       cx + Inches(0.06), y + Inches(0.025),
       _APP_COL_W - Inches(0.08), _CAT_HDR_H - Inches(0.03),
       size=8.5, bold=True, color=tc)
    y += _CAT_HDR_H

    # Task rows
    for row_i, (name, dur, tgt) in enumerate(tasks):
        bg = RGBColor(0xF5, 0xF5, 0xF5) if row_i % 2 == 0 else WHITE
        rect(sl, cx, y, _APP_COL_W, _TASK_H, fill=bg)
        star = " ★" if tgt else ""
        tb(sl, f"  {name}{star}",
           cx + Inches(0.04), y + Inches(0.010),
           _APP_COL_W * 0.74, _TASK_H - Inches(0.010),
           size=7.5, color=ORANGE if tgt else DGRAY, bold=bool(tgt))
        if dur and dur != '—':
            tb(sl, dur, cx + _APP_COL_W * 0.76, y + Inches(0.010),
               _APP_COL_W * 0.22, _TASK_H - Inches(0.010),
               size=7, color=MGRAY, align=PP_ALIGN.RIGHT)
        y += _TASK_H

    return y   # return updated y for next category in this column


# ── Appendix slide A1: All 65 Atomic Tasks ───────────────────────────────────
from robocasa.utils.dataset_registry import ATOMIC_TASK_DATASETS as _ATOMIC

sl = slide()
top = slide_heading(sl, "Appendix A — All 65 Atomic Tasks",
                    subtitle="★ = included in target evaluation split (18 tasks) · duration shown as seconds @ 20 Hz")

_ATOMIC_TARGET = {
    "CloseBlenderLid","CloseFridge","CloseToasterOvenDoor","CoffeeSetupMug",
    "NavigateKitchen","OpenCabinet","OpenDrawer","OpenStandMixerHead",
    "PickPlaceCounterToCabinet","PickPlaceCounterToStove","PickPlaceDrawerToCounter",
    "PickPlaceSinkToCounter","PickPlaceToasterToCounter","SlideDishwasherRack",
    "TurnOffStove","TurnOnElectricKettle","TurnOnMicrowave","TurnOnSinkFaucet",
}
_atomic_sorted = sorted(_ATOMIC.items())
_per_col = (_len := len(_atomic_sorted) + _COLS - 1) // _COLS   # ceil(65/3)=22
_per_col = 22   # fix to 22 so col 3 has 21

for ci in range(_COLS):
    cx = _col_x(ci)
    rect(sl, cx, top, _APP_COL_W, _CAT_HDR_H, fill=LBLUE_F, line_color=LBLUE, line_pt=0.5)
    lo = ci * _per_col + 1
    hi = min((ci + 1) * _per_col, 65)
    tb(sl, f"Atomic tasks  {lo}–{hi}",
       cx + Inches(0.06), top + Inches(0.025),
       _APP_COL_W - Inches(0.08), _CAT_HDR_H - Inches(0.03),
       size=8.5, bold=True, color=BLUE)

for idx, (task, info) in enumerate(_atomic_sorted):
    ci = idx // _per_col
    ri = idx % _per_col
    if ci >= _COLS:
        break
    cx = _col_x(ci)
    cy = top + _CAT_HDR_H + ri * _TASK_H
    bg = RGBColor(0xF5, 0xF5, 0xF5) if ri % 2 == 0 else WHITE
    rect(sl, cx, cy, _APP_COL_W, _TASK_H, fill=bg)
    h = info.get('horizon', 0) if isinstance(info, dict) else 0
    star = " ★" if task in _ATOMIC_TARGET else ""
    tb(sl, f"  {task}{star}",
       cx + Inches(0.04), cy + Inches(0.010),
       _APP_COL_W * 0.78, _TASK_H - Inches(0.010),
       size=7.5, color=ORANGE if star else DGRAY, bold=bool(star))
    if h:
        tb(sl, f"{h//20}s",
           cx + _APP_COL_W * 0.80, cy + Inches(0.010),
           _APP_COL_W * 0.18, _TASK_H - Inches(0.010),
           size=7, color=MGRAY, align=PP_ALIGN.RIGHT)

tb(sl, "★ = target evaluation split",
   Inches(0.5), top + _CAT_HDR_H + 22 * _TASK_H + Inches(0.05),
   Inches(5), Inches(0.25), size=8, color=MGRAY, italic=True)


# ── Appendix slides B1–BN: 60 Composite Categories ───────────────────────────
_cur_sl   = None
_cur_top  = None
_col_idx  = 0
_col_y    = Inches(0)
_page_num = 0
_page_cat = 0   # category colour-cycle counter per page


def _new_cat_slide():
    global _cur_sl, _cur_top, _col_idx, _col_y, _page_num, _page_cat
    _page_num += 1
    _cur_sl  = slide()
    _cur_top = slide_heading(
        _cur_sl,
        f"Appendix B — Composite Tasks by Category  (Part {_page_num})",
        subtitle="★ = target evaluation split  ·  duration shown at 20 Hz  ·  300 tasks across 60 activity categories")
    _col_idx = 0
    _col_y   = _cur_top
    _page_cat = 0


for cat_num, cat_name, tasks in _CATEGORIES:
    cat_h = _CAT_HDR_H + len(tasks) * _TASK_H + _CAT_GAP

    if _cur_sl is None:
        _new_cat_slide()

    # Does this category fit in the current column?
    if _col_y + cat_h > _cur_top + _MAX_COL_H:
        _col_idx += 1
        _col_y = _cur_top
        _page_cat += 1
        if _col_idx >= _COLS:
            _new_cat_slide()

    _draw_category(_cur_sl, _col_idx, _col_y, cat_num, cat_name, tasks, _page_cat)
    _col_y   += cat_h
    _page_cat += 1


# ── Legend note on last category slide ───────────────────────────────────────
if _cur_sl:
    tb(_cur_sl, "★ = included in composite-seen or composite-unseen target evaluation split",
       Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.28),
       size=8, color=MGRAY, italic=True)


# ── save ──────────────────────────────────────────────────────────────────────
out = "robocasa_slides.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
